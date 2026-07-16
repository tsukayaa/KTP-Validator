"""
Build the Stage-3 presentation deck (KTP screen-replay / pendar-bezel detector).

3 parts, same visual style as KTP_Validator_Overview.pptx:
  1) CLIP full fine-tune tested -> not strong enough (F1 ~74%)
  2) Forensic features used -> before/after panel (output/feature_panel.png)
  3) LightGBM result -> dataset composition + confusion matrix + metrics
     (reads output/results.json if present; falls back to a clearly-marked
     placeholder so the deck is always buildable, even before training finishes)

Run:
    python ktp_pad/viz_features.py          # produce output/feature_panel.png (di sebelah, real data)
    python ktp_pad/build_stage3_deck.py      # produce KTP_Stage3_Deck.pptx
"""
from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent.parent

SLIDE_W, SLIDE_H = Emu(12191695), Emu(6858000)

BG_DARK    = "1F2937"
BLUE       = "2563EB"
GRAY_LIGHT = "F3F4F6"
GRAY_MED   = "6B7280"
GRAY_SOFT  = "9CA3AF"
WARN_RED   = "F87171"
GREEN      = "16A34A"
GREEN_BG   = "DCFCE7"
AMBER      = "D97706"
AMBER_BG   = "FEF3C7"
RED        = "DC2626"
RED_BG     = "FEE2E2"
WHITE      = "FFFFFF"
TEXT_DARK  = "1F2937"

CLIP_F1 = 0.74


def rgb(h): return RGBColor.from_string(h)


def rect(slide, l, t, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(l), Emu(t), Emu(w), Emu(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = rgb(line); shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    """runs: list of (text, size_pt, bold, color_hex) OR list of paragraphs (list of runs)."""
    box = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    if runs and isinstance(runs[0], tuple):
        runs = [runs]
    for i, para_runs in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        for txt, size, bold, color in para_runs:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = rgb(color)
            r.font.name = "Calibri"
    return box


def title_bar(slide):
    rect(slide, 0, 0, SLIDE_W, 109728, BLUE)


def slide_header(slide, kicker_title, subtitle, subtitle_color=GRAY_MED):
    title_bar(slide)
    text(slide, 548640, 320040, 11064240, 640080,
         [(kicker_title, 28, True, TEXT_DARK)])
    text(slide, 548640, 932688, 11064240, 457200,
         [(subtitle, 14, False, subtitle_color)])


def footer(slide, page_num):
    text(slide, 548640, 6446520, 9144000, 320040,
         [("KTP Stage-3  •  Forensic features + LightGBM (pendar/bezel detector)", 10, False, GRAY_MED)])
    text(slide, 11155680, 6446520, 640080, 320040,
         [(str(page_num), 10, False, GRAY_MED)])


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 0 — cover
# ══════════════════════════════════════════════════════════════════════════
def build_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG_DARK)
    rect(slide, 0, 4160520, SLIDE_W, 54864, BLUE)
    text(slide, 822960, 2194560, 10515600, 1280160, [
        [("KTP Screen-Replay Detector", 44, True, WHITE)],
        [("Forensic features  +  LightGBM  —  kenapa CLIP tidak cukup", 20, False, "C7D2FE")],
    ])
    text(slide, 822960, 4343400, 10515600, 1097280, [
        [("Deteksi KTP asli vs KTP yang difoto ulang dari layar (“pendar” / bezel)", 16, False, GRAY_SOFT)],
    ])
    text(slide, 548640, 6446520, 9144000, 320040,
         [("KTP Stage-3  •  600/600 labeled, CPU, tanpa pretrained weights", 10, False, GRAY_MED)])
    text(slide, 11155680, 6446520, 640080, 320040, [("1", 10, False, GRAY_MED)])


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — CLIP: not strong enough
# ══════════════════════════════════════════════════════════════════════════
def build_clip_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "1 · CLIP full fine-tune: tidak cukup kuat",
                 "Sudah dicoba dulu — F1 mentok di sekitar 74%")

    box = rect(slide, 548640, 1691640, 5486400, 3612360, GRAY_LIGHT)
    text(slide, 822960, 1965960, 4937760, 3063240, [
        [("F1 score (test)", 14, True, GRAY_MED)],
        [("~74%", 48, True, AMBER)],
        [("", 8, False, TEXT_DARK)],
        [("plateau — tidak naik walau di-tuning lanjut", 12, False, GRAY_MED)],
    ])

    text(slide, 6400800, 1691640, 5212440, 365760,
         [("Kenapa CLIP structurally blind ke sinyal ini", 16, True, BLUE)])
    bullets = [
        "CLIP resize semua gambar ke 224×224",
        "CLIP dilatih supaya INVARIANT ke resize / JPEG recompress / color shift",
        "Tapi justru itu SIS sinyal yang dibutuhkan: moiré, blocking JPEG ganda, color banding",
        "Hasil: fitur high-frequency yang jadi bukti “layar” ikut hilang sebelum sempat dipelajari",
    ]
    runs = [[("•  " + b, 14, False, TEXT_DARK)] for b in bullets]
    text(slide, 6400800, 2148840, 5212440, 3200000, runs)

    text(slide, 548640, 5577840, 11064240, 640080, [
        [("Kesimpulan: butuh sinyal texture-level yang CLIP buang — bukan lagi fine-tune CLIP.", 14, True, TEXT_DARK)],
    ])
    footer(slide, 2)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — features used, before/after panel
# ══════════════════════════════════════════════════════════════════════════
def build_features_slide(prs, panel_png: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "2 · Fitur forensik yang dipakai",
                 "Texture/frequency signal yang CLIP buang, dihitung manual (~125 fitur)")

    if panel_png.exists():
        pic = slide.shapes.add_picture(str(panel_png), Emu(548640), Emu(1691640), width=Emu(7315200))
    else:
        rect(slide, 548640, 1691640, 7315200, 4300000, GRAY_LIGHT)
        text(slide, 822960, 3600000, 6766560, 640080,
             [("belum ada output/feature_panel.png — jalankan viz_features.py di sebelah", 13, False, GRAY_MED)],
             align=PP_ALIGN.CENTER)

    groups = [
        ("FFT (50 fitur)", "moiré: interferensi grid subpixel layar vs sensor kamera — paling kuat"),
        ("SRM residual (40 fitur)", "5 kernel high-pass — noise grid pixel layar vs grain kertas/PVC"),
        ("LBP (20 fitur)", "micro-texture: permukaan layar vs kertas/PVC"),
        ("Specular (5 fitur)", "blob terang low-saturation — “pendar”/glare layar"),
        ("Color (7 fitur)", "banding histogram — gamut warna layar terbatas"),
        ("Blockiness (3 fitur)", "diskontinuitas blok 8×8 — JPEG kompresi ganda"),
    ]
    x = 8100840
    y = 1691640
    for name, desc in groups:
        text(slide, x, y, 3542880, 640080, [
            [(name, 13, True, BLUE)],
            [(desc, 11, False, GRAY_MED)],
        ])
        y += 700000
    footer(slide, 3)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — dataset composition
# ══════════════════════════════════════════════════════════════════════════
def build_dataset_slide(prs, counts):
    train_pos, train_neg, test_pos, test_neg = counts
    total_all = train_pos + train_neg + test_pos + test_neg

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_header(slide, "3a · Komposisi dataset",
                 f"{train_pos+test_pos} asli (positives) / {train_neg+test_neg} spoof (negatives), labeled, no augmentation")

    tbl_shape = slide.shapes.add_table(3, 3, Emu(548640), Emu(1691640), Emu(5486400), Emu(1920240))
    tbl = tbl_shape.table
    tbl.columns[0].width = Emu(1828800)
    tbl.columns[1].width = Emu(1828800)
    tbl.columns[2].width = Emu(1828800)
    headers = ["", "Positives\n(asli)", "Negatives\n(spoof)"]
    rows = [["TRAIN", str(train_pos), str(train_neg)], ["TEST", str(test_pos), str(test_neg)]]

    def fmt_cell(cell, txt, bg, color, bold, size=13):
        cell.text = txt
        cell.fill.solid(); cell.fill.fore_color.rgb = rgb(bg)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = rgb(color)
                r.font.name = "Calibri"

    for c, h in enumerate(headers):
        fmt_cell(tbl.cell(0, c), h, BG_DARK, WHITE, True)
    for r, row in enumerate(rows, start=1):
        for c, v in enumerate(row):
            fmt_cell(tbl.cell(r, c), v, GRAY_LIGHT if c == 0 else WHITE, TEXT_DARK, c == 0, size=16)

    text(slide, 548640, 3977639, 7315200, 365760,
         [("Proporsi kelas per split (bar proporsional ke angka tabel di atas)", 14, True, BLUE)])

    def bar_row(y, label, pos, neg, total_w):
        text(slide, 548640, y, 1188720, 365760, [(label, 13, True, TEXT_DARK)])
        total = pos + neg
        w_pos = int(total_w * pos / total) if total else 0
        w_neg = total_w - w_pos
        rect(slide, 1828800, y + 18288, max(w_pos, 1), 384048, GREEN)
        rect(slide, 1828800 + w_pos, y + 18288, max(w_neg, 1), 384048, RED)
        if w_pos > 700000:
            text(slide, 1828800, y + 18288, w_pos, 384048, [(str(pos), 11, True, WHITE)],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if w_neg > 700000:
            text(slide, 1828800 + w_pos, y + 18288, w_neg, 384048, [(str(neg), 11, True, WHITE)],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    bar_row(4416552, "TRAIN", train_pos, train_neg, 6583680)
    bar_row(4965192, "TEST", test_pos, test_neg, 6583680)

    rect(slide, 1828800, 5577840, 164592, 164592, GREEN)
    text(slide, 2057400, 5541264, 1828800, 274320, [("positives (asli)", 12, False, GRAY_MED)])
    rect(slide, 4023360, 5577840, 164592, 164592, RED)
    text(slide, 4251960, 5541264, 1828800, 274320, [("negatives (spoof)", 12, False, GRAY_MED)])

    text(slide, 548640, 6100000, 11064240, 320040, [
        ("Catatan: 600 spoof adalah spoof yang sudah pernah tertangkap — F1 di sini adalah batas atas, bukan jaminan untuk spoof canggih yang belum pernah dilihat.", 11, False, GRAY_SOFT),
    ])
    footer(slide, 4)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — confusion matrix + metrics
# ══════════════════════════════════════════════════════════════════════════
def build_metrics_slide(prs, results: dict | None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    placeholder = results is None
    if placeholder:
        conf = {"TP": 90, "FP": 8, "TN": 92, "FN": 10}
        metrics = None
    else:
        conf = results["confusion"]
        metrics = results["metrics"]

    subtitle = "Kelas positif = “asli”"
    if placeholder:
        subtitle += "  ·  ANGKA PLACEHOLDER — jalankan train_pad.py lalu build ulang deck ini"
    slide_header(slide, "3b · Hasil metrik (test set)", subtitle,
                 subtitle_color=WARN_RED if placeholder else GRAY_MED)

    text(slide, 2423160, 1783080, 3474720, 274320, [("Prediksi →", 12, True, GRAY_MED)])
    text(slide, 137159, 3886200, 1280160, 274320, [("Aktual", 12, True, GRAY_MED)])
    text(slide, 1554480, 1783080, 1737360, 320040, [("ASLI (yes)", 12, True, TEXT_DARK)])
    text(slide, 3291839, 1783080, 1737360, 320040, [("SPOOF (no)", 12, True, TEXT_DARK)])
    text(slide, 0, 2788920, 1371600, 365760, [("Asli", 12, True, TEXT_DARK)])
    text(slide, 0, 4526280, 1371600, 365760, [("Spoof", 12, True, TEXT_DARK)])

    def cell(l, t, fillbg, label, val, color):
        rect(slide, l, t, 1737360, 1737360, fillbg)
        text(slide, l, t + 220000, 1737360, 1500000, [
            [(label, 13, True, color)],
            [(str(val), 30, True, color)],
        ], align=PP_ALIGN.CENTER)

    cell(1554480, 2148840, GREEN_BG, "TP", conf["TP"], GREEN)
    cell(3291839, 2148840, AMBER_BG, "FN", conf["FN"], AMBER)
    cell(1554480, 3886200, RED_BG, "FP", conf["FP"], RED)
    cell(3291839, 3886200, GREEN_BG, "TN", conf["TN"], GREEN)

    text(slide, 1554480, 5760720, 3474720, 548640, [
        ("FP = spoof lolos jadi “asli” (paling bahaya)", 11, False, RED),
    ])

    if metrics is not None:
        acc, prec, rec, f1 = metrics["accuracy"], metrics["precision_pos"], metrics["recall_pos"], metrics["f1_macro"]
    else:
        tp, fp, tn, fn = conf["TP"], conf["FP"], conf["TN"], conf["FN"]
        acc = (tp + tn) / (tp + fp + tn + fn)
        prec = tp / (tp + fp) if (tp + fp) else 0.0
        rec = tp / (tp + fn) if (tp + fn) else 0.0
        f1 = (2 * prec * rec / (prec + rec)) if (prec + rec) else 0.0

    def metric_box(l, t, label, value, color):
        rect(slide, l, t, 2880360, 1417320, WHITE, line=GRAY_LIGHT)
        text(slide, l, t + 137160, 2880360, 1200000, [
            [(label, 14, True, GRAY_MED)],
            [(f"{value*100:.2f}%", 34, True, color)],
        ], align=PP_ALIGN.CENTER)

    metric_box(5760720, 1965960, "Accuracy", acc, BLUE)
    metric_box(8961120, 1965960, "F1 (macro)", f1, GREEN)
    metric_box(5760720, 3703320, "Precision", prec, AMBER)
    metric_box(8961120, 3703320, "Recall", rec, RED)

    delta = f1 - CLIP_F1
    delta_color = GREEN if delta >= 0 else RED
    delta_txt = f"{'+' if delta>=0 else ''}{delta*100:.1f} pts vs CLIP baseline ({CLIP_F1*100:.0f}%)"
    text(slide, 5760720, 5257800, 5943600, 640080, [
        [(delta_txt, 16, True, delta_color)],
    ])
    text(slide, 5760720, 5850000, 5943600, 500000, [
        ("Accuracy = (TP+TN)/total  ·  Precision = TP/(TP+FP)  ·  Recall = TP/(TP+FN)  ·  F1 = 2·P·R/(P+R)", 10, False, GRAY_MED),
    ])
    footer(slide, 5)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--results", type=Path, default=ROOT / "output" / "results.json")
    ap.add_argument("--panel", type=Path, default=ROOT / "output" / "feature_panel.png")
    ap.add_argument("--out", type=Path, default=ROOT / "KTP_Stage3_Deck.pptx")
    ap.add_argument("--train-pos", type=int, default=500)
    ap.add_argument("--train-neg", type=int, default=500)
    ap.add_argument("--test-pos", type=int, default=100)
    ap.add_argument("--test-neg", type=int, default=100)
    args = ap.parse_args()

    results = None
    if args.results.exists():
        results = json.loads(args.results.read_text())
        print(f"using real results from {args.results}")
    else:
        print(f"WARNING: {args.results} not found -> slide 5 will use PLACEHOLDER numbers")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_cover(prs)
    build_clip_slide(prs)
    build_features_slide(prs, args.panel)
    build_dataset_slide(prs, (args.train_pos, args.train_neg, args.test_pos, args.test_neg))
    build_metrics_slide(prs, results)

    prs.save(str(args.out))
    print(f"saved: {args.out}")


if __name__ == "__main__":
    main()
